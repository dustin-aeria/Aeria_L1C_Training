import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def extract_pptx_content(pptx_path, output_dir):
    """Extract text and images from PowerPoint file"""

    # Create output directories
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)

    # Load presentation
    prs = Presentation(pptx_path)

    slides_data = []
    image_counter = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': [],
            'images': [],
            'notes': ''
        }

        # Extract text from shapes
        for shape in slide.shapes:
            # Get title
            if shape.has_text_frame:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # Title
                        slide_data['title'] = shape.text
                        continue

                # Get other text content
                if shape.text.strip():
                    slide_data['content'].append({
                        'type': 'text',
                        'text': shape.text
                    })

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_counter += 1
                image_filename = f"slide{slide_num:02d}_img{image_counter:02d}.{image.ext}"
                image_path = os.path.join(images_dir, image_filename)

                with open(image_path, 'wb') as f:
                    f.write(image.blob)

                slide_data['images'].append({
                    'filename': image_filename,
                    'path': f"/images/l1c00/{image_filename}"
                })

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                slide_data['notes'] = notes_frame.text

        slides_data.append(slide_data)

    # Save as JSON
    json_path = os.path.join(output_dir, 'slides_data.json')
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    print(f"Extracted {len(slides_data)} slides")
    print(f"Extracted {image_counter} images")
    print(f"Data saved to: {json_path}")

    return slides_data

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx.py <pptx_file> <output_dir>")
        sys.exit(1)

    pptx_file = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(pptx_file):
        print(f"Error: File not found: {pptx_file}")
        sys.exit(1)

    extract_pptx_content(pptx_file, output_dir)
